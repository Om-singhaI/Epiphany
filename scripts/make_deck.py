"""Generate the Epiphany pitch deck (docs/Epiphany_Pitch.pptx).

Run:  python scripts/make_deck.py
Produces a dark, on-brand 16:9 deck for the Google Cloud Rapid Agent Hackathon.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand palette ────────────────────────────────────────────────────────
BG = RGBColor(0x0B, 0x0F, 0x19)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
INDIGO = RGBColor(0x81, 0x8C, 0xF8)
CARD = RGBColor(0x15, 0x1B, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def accent_bar(s, x, y, w, color=INDIGO, h=0.07):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def para(tf, text, size, color=WHITE, bold=False, bullet=False, space=10, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = ("•  " + text) if bullet else text
    p.space_after = Pt(space)
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Calibri"
    return p


def title(s, kicker, head, color=INDIGO):
    accent_bar(s, 0.7, 0.7, 1.2, color)
    tf = box(s, 0.7, 0.85, 12, 1.4)
    para(tf, kicker, 14, color, bold=True, space=2, first=True)
    para(tf, head, 34, WHITE, bold=True, space=0)


def card(s, x, y, w, h, color=CARD):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.color.rgb = RGBColor(0x2A, 0x33, 0x48); shp.line.width = Pt(1)
    return shp


# ── 1. Title ───────────────────────────────────────────────────────────────
s = slide()
accent_bar(s, 0.7, 2.2, 2.0, EMERALD)
tf = box(s, 0.7, 2.4, 12, 3)
para(tf, "EPIPHANY", 60, WHITE, bold=True, space=4, first=True)
para(tf, "An Autonomous AI Data Scientist", 30, INDIGO, bold=True, space=18)
para(tf, "Point it at any dataset. It explores, hypothesizes, proves it with real "
        "statistics, trains a model, and ships it — on its own.", 18, SLATE, space=6)
tf2 = box(s, 0.7, 6.4, 12, 0.8)
para(tf2, "Google Cloud Rapid Agent Hackathon   •   live: epiphany-ds.fly.dev",
     14, EMERALD, bold=True, first=True)

# ── 2. Problem ──────────────────────────────────────────────────────────────
s = slide(); title(s, "THE PROBLEM", "Data sits unused. Data scientists don't scale.", EMERALD)
tf = box(s, 0.8, 2.4, 11.6, 4)
para(tf, "Most enterprise data is never analyzed — the questions outnumber the analysts.", 20, WHITE, bullet=True, first=True)
para(tf, "Hiring data scientists is slow and expensive; every question becomes a ticket in a queue.", 20, WHITE, bullet=True)
para(tf, "Today's “AI for data” tools mostly describe data or generate code you still have to run and trust.", 20, WHITE, bullet=True)
para(tf, "Nobody closes the loop: raw data → validated insight → a deployed model.", 20, EMERALD, bold=True, bullet=True)

# ── 3. Solution: the 5-step loop ────────────────────────────────────────────
s = slide(); title(s, "THE SOLUTION", "A continuous 5-step loop — and every step is real")
steps = [
    ("1 · TRIGGER", "Fivetran", "Wakes when new data syncs", INDIGO),
    ("2 · EXPLORE", "Elastic", "Discovers schema, ranks real signals", EMERALD),
    ("3 · REASON", "Gemini + ADK", "Forms a falsifiable hypothesis", INDIGO),
    ("4 · VALIDATE", "SciPy sandbox", "Runs the RIGHT test on real rows", EMERALD),
    ("5 · DEPLOY", "GitLab", "Trains a real model, opens an MR", INDIGO),
]
x = 0.7
for kick, tech, desc, col in steps:
    card(s, x, 2.6, 2.36, 3.2)
    tf = box(s, x + 0.18, 2.85, 2.0, 3.0)
    para(tf, kick, 13, col, bold=True, space=6, first=True)
    para(tf, tech, 17, WHITE, bold=True, space=8)
    para(tf, desc, 13, SLATE, space=0)
    x += 2.5
tf = box(s, 0.7, 6.1, 12, 0.8)
para(tf, "When Google ADK + Gemini are available, an LlmAgent dynamically chooses which "
        "tools to call. Otherwise the same real tools run as a direct pipeline.", 14, SLATE, first=True)

# ── 4. Why it's real (the differentiator) ────────────────────────────────────
s = slide(); title(s, "WHY IT WINS", "It's real data science — not a demo prop", EMERALD)
tf = box(s, 0.8, 2.3, 11.8, 4.6)
para(tf, "Real statistics.  SciPy tests on real rows — a weak signal returns NOT significant. The system can say “no.”", 18, WHITE, bullet=True, first=True)
para(tf, "Real models.  Trains a scikit-learn model, measures held-out ROC-AUC / R², saves a loadable .pkl artifact.", 18, WHITE, bullet=True)
para(tf, "Any dataset, any domain.  Auto-detects column types & target; picks χ² / t-test / ANOVA / correlation accordingly.", 18, WHITE, bullet=True)
para(tf, "Genuinely autonomous.  A background loop runs forever; the agent decides its own tool calls via Google ADK.", 18, WHITE, bullet=True)
para(tf, "Safe.  Agent-generated code is AST-screened and run in a network-isolated, resource-limited sandbox.", 18, WHITE, bullet=True)
para(tf, "Deployed & usable.  Live on the web, real auth (Clerk), bring-your-own-data, graceful degradation.", 18, EMERALD, bold=True, bullet=True)

# ── 5. Works on ANY dataset ──────────────────────────────────────────────────
s = slide(); title(s, "FOR ANYONE", "Upload any CSV — it adapts the test and the model")
rows = [
    ("Dataset", "Domain", "Target", "Test → Model", True),
    ("train.csv", "SaaS customers", "Churn (binary)", "t-test → classifier", False),
    ("wine_cultivars.csv", "Wine chemistry", "wine_class (3-class)", "ANOVA → classifier (94%)", False),
    ("diabetes_progression.csv", "Healthcare", "progression (numeric)", "correlation → regression", False),
]
y = 2.7
for c0, c1, c2, c3, head in rows:
    cols = [(c0, 3.3), (c1, 3.0), (c2, 3.2), (c3, 3.5)]
    x = 0.7
    for txt, w in cols:
        if head:
            card(s, x, y, w - 0.1, 0.6, RGBColor(0x1E, 0x29, 0x3B))
        tf = box(s, x + 0.12, y + 0.06, w - 0.2, 0.6)
        para(tf, txt, 14 if head else 14, INDIGO if head else WHITE, bold=head, first=True)
        x += w
    y += 0.72
tf = box(s, 0.7, 6.2, 12, 0.8)
para(tf, "Same agent. Zero configuration. It picks the statistically correct method for whatever you give it.", 15, EMERALD, bold=True, first=True)

# ── 6. Architecture / Google Cloud ───────────────────────────────────────────
s = slide(); title(s, "ARCHITECTURE", "Built on Google Cloud's agentic stack", EMERALD)
tf = box(s, 0.8, 2.3, 11.8, 4.4)
para(tf, "Gemini 2.5 Flash — the reasoning brain (API key or Vertex AI).", 18, WHITE, bullet=True, first=True)
para(tf, "Google Agent Development Kit (ADK) — orchestrates autonomous, dynamic tool-calling.", 18, WHITE, bullet=True)
para(tf, "FastAPI + WebSockets — live agent stream to a real-time dashboard.", 18, WHITE, bullet=True)
para(tf, "Hardened subprocess sandbox — AST scanner + no network egress + CPU/RAM limits.", 18, WHITE, bullet=True)
para(tf, "SQLite memory — every hypothesis & deployment persisted across cycles.", 18, WHITE, bullet=True)
para(tf, "Container-native — Dockerfile, honors $PORT, health checks: Cloud Run / Fly ready (deployed live).", 18, INDIGO, bold=True, bullet=True)

# ── 7. Live demo ─────────────────────────────────────────────────────────────
s = slide(); title(s, "LIVE DEMO", "See it think — in 90 seconds")
tf = box(s, 0.8, 2.3, 11.8, 4)
para(tf, "1.  Sign in (Clerk) → the dashboard's Active Agent Stream is already working autonomously.", 18, WHITE, bullet=True, first=True)
para(tf, "2.  Switch churn → wine: it re-profiles and now chooses ANOVA + a 94% classifier.", 18, WHITE, bullet=True)
para(tf, "3.  Run a mission in plain English → Gemini + ADK pick the tools live.", 18, WHITE, bullet=True)
para(tf, "4.  Upload your own CSV → it works on your data too.", 18, WHITE, bullet=True)
card(s, 0.8, 6.0, 11.7, 0.9, RGBColor(0x06, 0x2B, 0x22))
tf = box(s, 1.0, 6.18, 11.4, 0.7)
para(tf, "▶  Try it now:  https://epiphany-ds.fly.dev", 20, EMERALD, bold=True, first=True)

# ── 8. Close ─────────────────────────────────────────────────────────────────
s = slide()
accent_bar(s, 0.7, 2.3, 2.0, EMERALD)
tf = box(s, 0.7, 2.5, 12, 3)
para(tf, "Data science that runs itself.", 40, WHITE, bold=True, space=14, first=True)
para(tf, "Epiphany turns any dataset into a validated, deployed model — with no human "
        "in the loop. Real statistics, real models, any domain, live on the web.", 20, SLATE, space=18)
para(tf, "epiphany-ds.fly.dev", 22, INDIGO, bold=True, space=0)

out = Path("docs/Epiphany_Pitch.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"Saved {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
